"""
PowerPoint Slide Parser (Issue #9)

Processes PowerPoint files with proper slide-level structure preservation.
Groups content by slides in JSONL to maintain presentation context.

This module extracts:
- Slide-level grouping (each slide becomes a parent element)
- Slide titles and content
- Slide notes and speaker notes
- Slide numbers and hierarchy
- Images and shapes per slide
- Table data within slides

Example JSONL output:
    {
        "element_id": "slide_1",
        "type": "slide",
        "text": "Introduction to Cloud Migration",
        "metadata": {
            "slide_number": 1,
            "has_title": true,
            "has_notes": true,
            "shape_count": 5,
            "image_count": 2
        }
    }
    {
        "element_id": "slide_1_content_1",
        "type": "slide_content",
        "text": "Migration Strategy Overview...",
        "parent_id": "slide_1",
        "metadata": {"slide_number": 1, "content_type": "text_box"}
    }
"""
import logging
from typing import List, Dict, Any, Optional, Tuple
from hashlib import sha1
import os

logger = logging.getLogger(__name__)


class PowerPointSlideParser:
    """
    Parses PowerPoint presentations with slide-level structure preservation.
    
    Uses python-pptx library to extract slide-by-slide content while maintaining
    the hierarchical relationship between slides and their content.
    """
    
    def __init__(self):
        self.slide_elements = []
        self.stats = {}
    
    def parse_presentation(
        self,
        file_path: str,
        filename: str
    ) -> Tuple[List[Dict[str, Any]], Dict[str, Any]]:
        """
        Parse PowerPoint file into structured elements.
        
        Args:
            file_path: Path to PPTX file
            filename: Original filename
        
        Returns:
            Tuple of (elements, stats)
        """
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx not available for PowerPoint parsing")
            return [], {"error": "python-pptx not installed"}
        
        try:
            prs = Presentation(file_path)
            elements = []
            
            total_slides = len(prs.slides)
            slide_with_notes = 0
            total_shapes = 0
            total_images = 0
            total_tables = 0
            
            for slide_idx, slide in enumerate(prs.slides, start=1):
                # Create slide parent element
                slide_element, slide_stats = self._process_slide(
                    slide=slide,
                    slide_number=slide_idx,
                    filename=filename
                )
                elements.append(slide_element)
                
                # Process slide content (shapes, text boxes, images, tables)
                content_elements = self._process_slide_content(
                    slide=slide,
                    slide_number=slide_idx,
                    parent_id=slide_element["element_id"],
                    filename=filename
                )
                elements.extend(content_elements)
                
                # Process slide notes if present
                if slide.has_notes_slide:
                    slide_with_notes += 1
                    notes_element = self._process_slide_notes(
                        slide=slide,
                        slide_number=slide_idx,
                        parent_id=slide_element["element_id"],
                        filename=filename
                    )
                    if notes_element:
                        elements.append(notes_element)
                
                # Update stats
                total_shapes += slide_stats.get("shape_count", 0)
                total_images += slide_stats.get("image_count", 0)
                total_tables += slide_stats.get("table_count", 0)
            
            # Build stats
            stats = {
                "total_slides": total_slides,
                "slides_with_notes": slide_with_notes,
                "total_shapes": total_shapes,
                "total_images": total_images,
                "total_tables": total_tables,
                "total_elements": len(elements),
                "parser": "python-pptx",
                "format": "pptx"
            }
            
            logger.info(
                f"PowerPoint parsing complete: {total_slides} slides, "
                f"{len(elements)} elements extracted from '{filename}'"
            )
            
            return elements, stats
            
        except Exception as e:
            logger.error(f"PowerPoint parsing failed for '{filename}': {e}", exc_info=True)
            return [], {"error": str(e)}
    
    def _process_slide(
        self,
        slide: Any,
        slide_number: int,
        filename: str
    ) -> Tuple[Dict[str, Any], Dict[str, Any]]:
        """Process a single slide and create parent element."""
        # Extract slide title
        title_text = self._extract_slide_title(slide)
        
        # Count different content types
        shape_count = 0
        image_count = 0
        table_count = 0
        text_box_count = 0
        
        for shape in slide.shapes:
            shape_count += 1
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_count += 1
            elif hasattr(shape, "has_table") and shape.has_table:
                table_count += 1
            elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
                text_box_count += 1
        
        # Create element ID
        element_id = self._generate_element_id(filename, f"slide_{slide_number}")
        
        # Build slide element
        slide_element = {
            "element_id": element_id,
            "type": "slide",
            "text": title_text or f"Slide {slide_number}",
            "page_number": slide_number,
            "coordinates": None,
            "parent_id": None,
            "metadata": {
                "slide_number": slide_number,
                "has_title": bool(title_text),
                "has_notes": slide.has_notes_slide,
                "shape_count": shape_count,
                "image_count": image_count,
                "table_count": table_count,
                "text_box_count": text_box_count,
                "source": "powerpoint_slide",
                "filename": filename
            },
            "hierarchy_level": 0,
            "semantic_tags": ["presentation_slide", "slide_header"],
            "confidence_score": 1.0
        }
        
        stats = {
            "shape_count": shape_count,
            "image_count": image_count,
            "table_count": table_count
        }
        
        return slide_element, stats
    
    def _process_slide_content(
        self,
        slide: Any,
        slide_number: int,
        parent_id: str,
        filename: str
    ) -> List[Dict[str, Any]]:
        """Process all content shapes within a slide."""
        content_elements = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            # Skip title shape (already extracted)
            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                placeholder = shape.placeholder_format
                if placeholder.type == 1:  # PP_PLACEHOLDER.TITLE
                    continue
            
            # Process based on shape type
            if hasattr(shape, "has_table") and shape.has_table:
                # Table shape
                table_element = self._process_table_shape(
                    shape=shape,
                    slide_number=slide_number,
                    shape_index=shape_idx,
                    parent_id=parent_id,
                    filename=filename
                )
                if table_element:
                    content_elements.append(table_element)
            
            elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
                # Text box or shape with text
                text_element = self._process_text_shape(
                    shape=shape,
                    slide_number=slide_number,
                    shape_index=shape_idx,
                    parent_id=parent_id,
                    filename=filename
                )
                if text_element:
                    content_elements.append(text_element)
            
            elif shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                # Image shape
                image_element = self._process_image_shape(
                    shape=shape,
                    slide_number=slide_number,
                    shape_index=shape_idx,
                    parent_id=parent_id,
                    filename=filename
                )
                if image_element:
                    content_elements.append(image_element)
        
        return content_elements
    
    def _process_text_shape(
        self,
        shape: Any,
        slide_number: int,
        shape_index: int,
        parent_id: str,
        filename: str
    ) -> Optional[Dict[str, Any]]:
        """Extract text from a text shape."""
        try:
            text = shape.text.strip()
            if not text:
                return None
            
            element_id = self._generate_element_id(
                filename,
                f"slide_{slide_number}_text_{shape_index}"
            )
            
            return {
                "element_id": element_id,
                "type": "slide_content",
                "text": text,
                "page_number": slide_number,
                "coordinates": None,
                "parent_id": parent_id,
                "metadata": {
                    "slide_number": slide_number,
                    "content_type": "text_box",
                    "shape_index": shape_index,
                    "source": "powerpoint_shape",
                    "filename": filename
                },
                "hierarchy_level": 1,
                "semantic_tags": ["presentation_content", "text"],
                "confidence_score": 1.0
            }
        except Exception as e:
            logger.warning(f"Failed to extract text from shape {shape_index}: {e}")
            return None
    
    def _process_table_shape(
        self,
        shape: Any,
        slide_number: int,
        shape_index: int,
        parent_id: str,
        filename: str
    ) -> Optional[Dict[str, Any]]:
        """Extract table data from a table shape."""
        try:
            table = shape.table
            rows = table.rows
            
            # Extract table as text representation
            table_data = []
            for row in rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            
            # Create readable table representation
            table_text = "\n".join([" | ".join(row) for row in table_data])
            
            element_id = self._generate_element_id(
                filename,
                f"slide_{slide_number}_table_{shape_index}"
            )
            
            return {
                "element_id": element_id,
                "type": "slide_table",
                "text": table_text,
                "page_number": slide_number,
                "coordinates": None,
                "parent_id": parent_id,
                "metadata": {
                    "slide_number": slide_number,
                    "content_type": "table",
                    "shape_index": shape_index,
                    "table_rows": len(table_data),
                    "table_cols": len(table_data[0]) if table_data else 0,
                    "table_data": table_data,
                    "source": "powerpoint_table",
                    "filename": filename
                },
                "hierarchy_level": 1,
                "semantic_tags": ["presentation_content", "table", "structured_data"],
                "confidence_score": 1.0
            }
        except Exception as e:
            logger.warning(f"Failed to extract table from shape {shape_index}: {e}")
            return None
    
    def _process_image_shape(
        self,
        shape: Any,
        slide_number: int,
        shape_index: int,
        parent_id: str,
        filename: str
    ) -> Optional[Dict[str, Any]]:
        """Extract image metadata from an image shape."""
        try:
            element_id = self._generate_element_id(
                filename,
                f"slide_{slide_number}_image_{shape_index}"
            )
            
            # Try to get image name/description
            image_name = getattr(shape, "name", f"Image {shape_index}")
            
            return {
                "element_id": element_id,
                "type": "slide_image",
                "text": f"[Image: {image_name}]",
                "page_number": slide_number,
                "coordinates": None,
                "parent_id": parent_id,
                "metadata": {
                    "slide_number": slide_number,
                    "content_type": "image",
                    "shape_index": shape_index,
                    "image_name": image_name,
                    "source": "powerpoint_image",
                    "filename": filename
                },
                "hierarchy_level": 1,
                "semantic_tags": ["presentation_content", "image"],
                "confidence_score": 1.0
            }
        except Exception as e:
            logger.warning(f"Failed to extract image from shape {shape_index}: {e}")
            return None
    
    def _process_slide_notes(
        self,
        slide: Any,
        slide_number: int,
        parent_id: str,
        filename: str
    ) -> Optional[Dict[str, Any]]:
        """Extract speaker notes from a slide."""
        try:
            if not slide.has_notes_slide:
                return None
            
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text = notes_text_frame.text.strip()
            
            if not notes_text:
                return None
            
            element_id = self._generate_element_id(
                filename,
                f"slide_{slide_number}_notes"
            )
            
            return {
                "element_id": element_id,
                "type": "slide_notes",
                "text": notes_text,
                "page_number": slide_number,
                "coordinates": None,
                "parent_id": parent_id,
                "metadata": {
                    "slide_number": slide_number,
                    "content_type": "speaker_notes",
                    "source": "powerpoint_notes",
                    "filename": filename
                },
                "hierarchy_level": 1,
                "semantic_tags": ["presentation_notes", "speaker_notes"],
                "confidence_score": 1.0
            }
        except Exception as e:
            logger.warning(f"Failed to extract notes from slide {slide_number}: {e}")
            return None
    
    def _extract_slide_title(self, slide: Any) -> Optional[str]:
        """Extract title from a slide."""
        try:
            if slide.shapes.title:
                return slide.shapes.title.text.strip()
        except:
            pass
        
        # Fallback: Look for title placeholder
        for shape in slide.shapes:
            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                placeholder = shape.placeholder_format
                if placeholder.type == 1:  # PP_PLACEHOLDER.TITLE
                    if hasattr(shape, "text_frame"):
                        return shape.text_frame.text.strip()
        
        return None
    
    def _generate_element_id(self, filename: str, identifier: str) -> str:
        """Generate stable element ID."""
        combined = f"{filename}|{identifier}".encode('utf-8', 'ignore')
        return sha1(combined).hexdigest()


def parse_powerpoint_file(
    file_path: str,
    filename: str
) -> Tuple[List[Dict[str, Any]], Dict[str, Any]]:
    """
    Convenience function to parse PowerPoint file.
    
    Args:
        file_path: Path to PPTX file
        filename: Original filename
    
    Returns:
        Tuple of (elements, stats)
    """
    parser = PowerPointSlideParser()
    return parser.parse_presentation(file_path, filename)
